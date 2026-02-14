from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Data Structure: (Title, [Bullet Points], Speaker Notes)
    slides_data = [
        # Slide 1
        ("In Silico Approaches for Novel Insecticidal Target Identification",
         ["MSc Dissertation Proposal & Tool Demonstration",
          "Domain: Bioinformatics & Chemoinformatics",
          "Candidate: MSc Student"],
         "Good morning respected HOD and faculty. Today I am presenting the progress of my MSc dissertation, which focuses on building a computational suite to discover and model novel protein targets for agricultural pest management."),
        
        # Slide 2
        ("1. Introduction & Background",
         ["Traditional pest control relies heavily on broad-spectrum chemical insecticides.",
          "This leads to two major crises:",
          "  - Environmental toxicity and harm to non-target species.",
          "  - Rapid development of pesticide resistance in insect populations.",
          "The modern alternative is 'Rational Pesticide Design': targeting species-specific proteins."],
         "To give some background, modern agriculture is fighting a losing battle against pesticide resistance. Furthermore, the chemicals we use are often toxic to bees and aquatic life. The goal of modern research is to design chemicals that only target specific proteins unique to the pest, leaving other organisms unharmed."),
        
        # Slide 3
        ("2. The Problem Statement",
         ["Rational design requires structural data, but we face a 'Data Scarcity' bottleneck.",
          "Non-Model Organisms: Most agricultural pests do not have fully sequenced genomes.",
          "Orphan Ligands: We often discover chemicals that kill insects, but we don't know WHICH protein they bind to.",
          "In Vitro Cost: Identifying targets and folding proteins in a wet lab takes years and massive funding."],
         "The main problem we face is a massive data gap. Pests like the Cotton Bollworm aren't studied as much as humans or mice. We often find chemical compounds that work, but we don't know their protein target. And even if we guess the target, we don't have its 3D structure for that specific pest species to study how the drug binds."),
        
        # Slide 4
        ("3. Project Objectives",
         ["Develop an integrated, web-based bioinformatics suite to bypass the data scarcity bottleneck.",
          "Objective 1: Predict probable protein targets for novel chemical compounds using cheminformatics.",
          "Objective 2: Automate the discovery of orthologous target sequences in unsequenced pest transcriptomes.",
          "Objective 3: Predict the 3D tertiary structure of these targets in real-time.",
          "Objective 4: Ensure the tool is accessible via a user-friendly Streamlit dashboard."],
         "Therefore, the objective of my dissertation is to build a two-module software suite. First, a tool to predict what protein a chemical will attack. Second, a tool to find that protein's sequence in a pest's DNA and predict its 3D shape, all packaged in a clean web interface."),
        
        # Slide 5
        ("4. Module 1: Ligand-Based Target Prediction",
         ["The 'Chemoinformatics Engine' of the suite.",
          "Core Principle: The Similar Property Principle.",
          "  - Structurally similar molecules tend to exhibit similar biological activities and bind to similar targets.",
          "Function: Takes a user's chemical SMILES string and predicts which insect protein it will disrupt."],
         "Let's look at Module 1. This is the chemoinformatics engine. It operates on a fundamental law of medicinal chemistry: if two chemicals look similar, they probably do the same thing. If a user inputs a new chemical, this tool predicts its target based on what similar chemicals do."),
        
        # Slide 6
        ("5. Module 1: Datasets & Curation",
         ["Powered by two distinct datasets:",
          "1. ChEMBL Database Extraction:",
          "  - Filtered specifically for insect bioassays and bioactive compounds.",
          "2. Literature-Curated Dataset:",
          "  - High-confidence data extracted manually from published literature.",
          "  - Strong focus on established targets like Acetylcholinesterase (AChE)."],
         "For this to work, I needed high-quality data. I procured a dataset from the ChEMBL database, specifically filtering for insect targets. To make it highly contestable for my viva, I also integrated a secondary dataset manually curated from literature to provide a 'ground truth' baseline."),
        
        # Slide 7
        ("6. Module 1: Algorithmic Methodology",
         ["Powered by RDKit (Open-source cheminformatics software).",
          "Feature Extraction: Morgan Fingerprints.",
          "  - Converts 2D molecular structures into 2048-bit binary vectors (Radius 2).",
          "  - Captures the exact atomic neighborhoods of the molecule.",
          "Similarity Metric: Tanimoto Coefficient.",
          "  - Calculates the mathematical overlap between the query's fingerprint and the database compounds."],
         "Under the hood, Module 1 uses RDKit. It takes the chemical structure and converts it into a Morgan Fingerprint—a mathematical representation of its atoms. It then uses the Tanimoto coefficient to mathematically score how closely this new chemical matches thousands of known insect poisons in our database."),
        
        # Slide 8
        ("7. Module 1: Scoring & Validation",
         ["Mean Similarity Scoring:",
          "  - Aggregates the Tanimoto scores of all matched ligands for a specific target group.",
          "Confidence Labeling:",
          "  - Dynamically labels predictions as 'High', 'Medium', or 'Low' based on cluster size.",
          "The 'Literature Bonus':",
          "  - If a high-similarity match (>0.6) is found in the curated literature dataset, the algorithm applies a weighted bonus to the final ChEMBL score."],
         "To ensure accuracy, the tool doesn't just look at one match. It groups matches by target and calculates a mean score. If the chemical matches a highly verified target from my literature dataset, the algorithm applies a statistical 'Literature Bonus' to boost its rank, ensuring we prioritize validated biological pathways."),
        
        # Slide 9
        ("8. Module 2: Homology Modeler (The Bridge)",
         ["Once a target is predicted (e.g., Target X), we face a new problem:",
          "  - We know Target X exists in Fruit Flies, but what does it look like in the Diamondback Moth?",
          "Module 2 is an 'In Silico Gene Discovery' pipeline.",
          "It bridges the gap between known sequences and raw pest transcriptomes."],
         "Now we transition to Module 2. Let's say Module 1 predicts our chemical attacks 'Protein X'. We might have the sequence for Protein X in a common fruit fly, but not in the specific pest destroying crops. Module 2 solves this by finding the equivalent protein in the pest's raw DNA."),
        
        # Slide 10
        ("9. Module 2: Transcriptome Mining",
         ["Direct Integration with NCBI Remote APIs.",
          "Algorithms utilized:",
          "  - TBLASTN: Compares a protein query against a translated nucleotide database.",
          "  - BLASTN: Compares a nucleotide query against a nucleotide database.",
          "Allows researchers to search massive, unannotated genomic data dynamically without local supercomputers."],
         "Module 2 connects directly to NCBI's supercomputers via an API. It takes our reference protein and uses the 'tblastn' algorithm to scan the entire raw nucleotide database of the pest species, hunting for genetic fragments that match our target."),
        
        # Slide 11
        ("10. Module 2: Sequence Assembly & ORF Finder",
         ["BLAST returns fragmented matches (High-Scoring Segment Pairs).",
          "The tool automatically cleans and assembles these fragments.",
          "Open Reading Frame (ORF) Extraction:",
          "  - Custom algorithm translates the raw DNA across all 3 forward reading frames.",
          "  - Scans for start codons (Methionine) and stop codons.",
          "  - Extracts the longest, most viable functional protein sequence."],
         "When NCBI finds a match, it usually returns raw, fragmented DNA. My Python script acts as an assembler. It cleans the gaps, translates the DNA, and runs a custom Open Reading Frame algorithm to find the exact start and stop points of the functional protein, effectively doing 'in silico cloning'."),
        
        # Slide 12
        ("11. Module 2: 3D Structure Prediction",
         ["Traditional methods (AlphaFold, Swiss-Model) are too slow for real-time web tools.",
          "Solution: Integration of the ESMFold API (Meta AI).",
          "  - Uses a massive protein language model (LLM for biology).",
          "  - Folds the newly discovered target sequence into a 3D PDB structure in milliseconds.",
          "  - Rendered interactively in the browser using py3Dmol."],
         "Once we have the pest's exact protein sequence, we need its 3D shape to see how drugs bind to it. Instead of waiting hours for AlphaFold, I integrated Meta's ESMFold API. It uses AI language models to accurately predict the 3D structure in milliseconds, which is then rendered directly on the screen."),
        
        # Slide 13
        ("12. The Synergistic Workflow",
         ["How the suite functions end-to-end for a researcher:",
          "1. Input novel molecule.",
          "2. Module 1 identifies 'Acetylcholinesterase' as the likely target.",
          "3. User pushes the generic AChE sequence to Module 2.",
          "4. Module 2 isolates the specific AChE gene in the target pest.",
          "5. Module 2 folds the pest-specific AChE 3D structure.",
          "Result: The researcher is immediately ready for molecular docking simulations."],
         "Here is the workflow in action: A chemist creates a new molecule. Module 1 tells them it likely targets AChE. They take a standard AChE sequence, plug it into Module 2, type in their specific pest, and within two minutes, they have a 3D model of the pest's AChE ready for docking analysis."),
        
        # Slide 14
        ("13. Technology Stack",
         ["Front-End / UI:",
          "  - Streamlit (Python-based reactive web framework)",
          "Back-End / Bioinformatics:",
          "  - BioPython (NCBI WWW, Seq, Entrez)",
          "Chemoinformatics:",
          "  - RDKit, Pandas",
          "Visualization & APIs:",
          "  - py3Dmol, ESMFold REST API"],
         "This entire suite was built using Python. Streamlit handles the web dashboard, allowing for a clean, multipage app. RDKit handles the heavy chemical math, BioPython manages the complex biological sequence parsing and NCBI communication, and Py3DMol handles the 3D graphics."),
        
        # Slide 15
        ("14. Phase 3: Mentor's Addition - Hepatotoxicity",
         ["Future Work & Final Implementation.",
          "Problem: Many great insecticidal compounds fail because they cause Drug-Induced Liver Injury (DILI) in humans/mammals.",
          "Proposed Addition: A Hepatotoxicity Predictor.",
          "Will utilize a similar RDKit/Machine Learning approach to flag compounds that share structural motifs with known liver toxins.",
          "Ensures early-stage safety screening."],
         "Following discussions with my mentor, we identified a missing piece: safety. An insecticide is useless if it damages human livers. Before submitting my thesis, I will be adding a third module: a Hepatotoxicity Predictor. This will flag compounds that might be dangerous to mammals early in the pipeline."),
        
        # Slide 16
        ("15. Significance & Impact",
         ["Bypasses the 'Orphan Target' problem in agrochemicals.",
          "Democratizes advanced bioinformatics:",
          "  - Replaces complex command-line pipelines with a clean graphical interface.",
          "Reduces Time & Cost:",
          "  - Compresses months of wet-lab target identification and cloning into a few minutes of computational time."],
         "The ultimate significance of this dissertation is accessibility. It takes complex chemoinformatics and command-line genomic mining and puts it behind a simple web interface. It compresses what would normally take months of wet-lab work into a pipeline that runs in under five minutes."),
        
        # Slide 17
        ("16. Conclusion",
         ["Modules 1 and 2 are successfully developed, integrated, and functioning.",
          "The pipeline successfully links chemical queries to structural biology.",
          "Next Steps:",
          "  - Finalize the Hepatotoxicity module.",
          "  - Complete rigorous edge-case testing.",
          "  - Begin drafting the thesis manuscript.",
          "Thank you. Open for feedback from the review committee."],
         "In conclusion, the core framework of the suite is fully operational. The integration of chemical similarity and structural homology works seamlessly. My next steps are finishing the toxicity module and writing the thesis. Thank you for your time, I am happy to take any questions or feedback from the committee.")
    ]

    for i, (title, points, notes) in enumerate(slides_data):
        # Determine layout (Title slide for first, Bullet slide for rest)
        if i == 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            for point in points:
                p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(20)
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Bullets
            tf = slide.placeholders[1].text_frame
            for point in points:
                p = tf.add_paragraph()
                p.text = point
                # Adjust font size for better fit
                p.font.size = Pt(22)
                # Indent sub-bullets
                if point.startswith("  -"):
                    p.level = 1
                    p.font.size = Pt(18)
        
        # ADD SPEAKER NOTES
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

    # Save the presentation
    filename = "HOD_Dissertation_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Success! Your detailed presentation has been saved as: {filename}")
    print("Speaker notes have been automatically embedded into every slide.")

if __name__ == '__main__':
    create_presentation()
